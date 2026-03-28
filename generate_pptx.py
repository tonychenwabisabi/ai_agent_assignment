from pptx import Presentation
from pptx.util import Inches
import os

# initialize
prs = Presentation()

def add_slide(title, bullets):
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)

    title_placeholder = slide.shapes.title
    content = slide.placeholders[1]

    title_placeholder.text = title

    tf = content.text_frame
    tf.clear()

    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

# read markdown files
def parse_md(file_path):
    slides = []
    with open(file_path, "r") as f:
        lines = f.readlines()

    current_title = None
    bullets = []

    for line in lines:
        line = line.strip()

        if line.startswith("# "):  # 新 slide
            if current_title:
                slides.append((current_title, bullets))
            current_title = line.replace("# ", "")
            bullets = []

        elif line.startswith("- "):
            bullets.append(line.replace("- ", ""))

    if current_title:
        slides.append((current_title, bullets))

    return slides

# add title slide
add_slide("AI Agent Assignment", ["Multi-tenant scaling", "Claude Code Efficiency", "Cognitive limits"])

# read three markdown files
files = [
    ("Q1: Multi-tenant Scaling", "outputs/multi-tenant-scaling-analysis.md"),
    ("Q2: Claude Code Efficiency", "outputs/double-productivity-claude-code.md"),
    ("Q3: Cognitive Limits", "outputs/cognitive-limits-analysis.md")
]

for section_title, file in files:
    # section slide
    add_slide(section_title, [])

    slides = parse_md(file)
    for title, bullets in slides:
        add_slide(title, bullets[:10])  # restrict the bullet number in one slide to be no more than 10.

# save pptx
prs.save("presentation.pptx")

print("✅ PPT generated: presentation.pptx")